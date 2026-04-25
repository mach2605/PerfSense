"""
Generate Mid-Semester Viva Presentation for PerfSense Dissertation
Matches theme of docs/PerfSense_Viva_Presentation_abstract_viva.pptx
Output: docs/reports/MS2024TM93208_MidSem_Viva.pptx
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette (matches existing PPT) ──────────────────────────────────────
BLUE_DARK   = RGBColor(0x1F, 0x47, 0x88)   # #1F4788  titles
BLUE_MID    = RGBColor(0x2E, 0x75, 0xB6)   # #2E75B6  sub-headings / accents
GREEN_OK    = RGBColor(0x37, 0x86, 0x36)   # #378636  good results
RED_WARN    = RGBColor(0xC0, 0x00, 0x00)   # #C00000  warnings / challenges
BLACK       = RGBColor(0x00, 0x00, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF2, 0xF2)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.63)

# ── Helper: create blank presentation ─────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Low-level helpers ──────────────────────────────────────────────────────────
def add_textbox(slide, left, top, width, height, text,
                bold=False, size=20, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    return txBox


def add_textbox_lines(slide, left, top, width, height, lines,
                      default_size=16, default_color=BLACK,
                      default_bold=False, wrap=True, line_spacing=None):
    """
    lines: list of (text, bold, size, color, align)
           or just str  → uses defaults with bullet=False
           or dict with any subset of keys
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap

    first = True
    for line in lines:
        if isinstance(line, str):
            spec = {"text": line}
        else:
            spec = line

        text   = spec.get("text", "")
        bold   = spec.get("bold",  default_bold)
        size   = spec.get("size",  default_size)
        color  = spec.get("color", default_color)
        align  = spec.get("align", PP_ALIGN.LEFT)
        indent = spec.get("indent", 0)   # indent level (0=none, 1=bullet)

        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()

        para.alignment = align
        if indent:
            para.level = indent

        run = para.add_run()
        run.text = text
        run.font.bold  = bold
        run.font.size  = Pt(size)
        run.font.color.rgb = color

    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def slide_header(slide, title, subtitle=None):
    """Standard dark-blue bold title at top."""
    if subtitle:
        # Compact layout so the full header fits within ~0.78" — content at y≥0.88 stays clear
        add_textbox(slide, 0.3, 0.05, 9.5, 0.50, title,
                    bold=True, size=26, color=BLUE_DARK)
        add_rect(slide, 0.3, 0.57, 9.4, 0.03, BLUE_DARK)
        add_textbox(slide, 0.3, 0.61, 9.4, 0.20, subtitle,
                    bold=False, size=12, color=BLUE_MID)
        # subtitle ends at 0.81" → gap of ≥0.07" before content at 0.88"
    else:
        add_textbox(slide, 0.3, 0.1, 9.5, 0.6, title,
                    bold=True, size=28, color=BLUE_DARK)
        add_rect(slide, 0.3, 0.72, 9.4, 0.04, BLUE_DARK)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE

add_textbox(slide, 0.5, 0.8, 9.0, 1.1,
            "PerfSense: ML-Based Performance Regression\nPredictor for React Applications",
            bold=True, size=30, color=BLUE_DARK, align=PP_ALIGN.CENTER)

add_rect(slide, 1.5, 2.05, 7.0, 0.05, BLUE_MID)

add_textbox_lines(slide, 0.5, 2.2, 9.0, 2.8,
    [
        {"text": "Mid-Semester Viva Presentation", "bold": True,  "size": 18, "color": BLUE_MID,  "align": PP_ALIGN.CENTER},
        {"text": "",                                "size": 10},
        {"text": "CHETHAN S",                       "bold": True,  "size": 16, "color": BLACK,      "align": PP_ALIGN.CENTER},
        {"text": "BITS ID: 2024TM93208",             "size": 14,   "color": BLACK,                   "align": PP_ALIGN.CENTER},
        {"text": "MTech in Software Engineering | Specialization: Full Stack Engineering", "size": 13, "color": BLACK, "align": PP_ALIGN.CENTER},
        {"text": "Course: S2-25_SEZG628T (Dissertation)",  "size": 13, "color": BLACK,               "align": PP_ALIGN.CENTER},
        {"text": "Supervisor: Ankit Kagliwal",       "size": 13,   "color": BLACK,                   "align": PP_ALIGN.CENTER},
        {"text": "Cadence Design Systems, Bengaluru","size": 13,   "color": BLACK,                   "align": PP_ALIGN.CENTER},
        {"text": "",                                 "size": 10},
        {"text": "April 2026",                       "bold": True,  "size": 14, "color": BLUE_MID,  "align": PP_ALIGN.CENTER},
    ])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Presentation Outline
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Presentation Outline")

items = [
    ("1.  Introduction & Problem Statement",  BLUE_DARK,  True),
    ("2.  Literature Review",                 BLACK,      False),
    ("3.  Research Methodology",              BLACK,      False),
    ("4.  Data Collection & Feature Engineering", BLACK,  False),
    ("5.  ML Model Results",                  BLACK,      False),
    ("6.  Case Studies – Real-World Validation", BLACK,   False),
    ("7.  Implementation Status & Next Steps",BLACK,      False),
    ("8.  Challenges & Mitigation",           BLACK,      False),
    ("9.  Summary",                           BLACK,      False),
]
lines = []
for txt, col, bold in items:
    lines.append({"text": txt, "bold": bold, "size": 17, "color": col})

add_textbox_lines(slide, 1.2, 0.95, 7.8, 4.3, lines)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Introduction: Problem Statement
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Introduction: Problem Statement")

left_lines = [
    {"text": "Performance Issues Are Expensive", "bold": True, "size": 16, "color": BLUE_DARK},
    {"text": "• 1-second delay → 7% drop in conversions", "size": 14, "color": BLACK},
    {"text": "• 1-second delay → 11% fewer page views", "size": 14, "color": BLACK},
    {"text": "• React apps: complexity ↑ = regression risk ↑", "size": 14, "color": BLACK},
    {"text": "", "size": 8},
    {"text": "Current Detection is Reactive", "bold": True, "size": 16, "color": BLUE_DARK},
    {"text": "• Issues discovered AFTER production deployment", "size": 14, "color": BLACK},
    {"text": "• Lighthouse, New Relic, Datadog: post-deploy only", "size": 14, "color": BLACK},
    {"text": "• Manual code reviews: 4–8 hrs/incident to debug", "size": 14, "color": BLACK},
    {"text": "", "size": 8},
    {"text": "The Gap", "bold": True, "size": 16, "color": RED_WARN},
    {"text": "No tool predicts performance impact BEFORE merge", "bold": True, "size": 14, "color": RED_WARN},
]
add_textbox_lines(slide, 0.4, 0.9, 5.6, 4.4, left_lines)

# Right box: highlight
add_rect(slide, 6.2, 0.9, 3.5, 4.3, LIGHT_GREY, BLUE_MID)
add_textbox_lines(slide, 6.4, 1.0, 3.1, 4.1,
    [
        {"text": "What Existing Tools Do", "bold": True, "size": 13, "color": BLUE_DARK},
        {"text": "✓ Collect runtime metrics", "size": 12, "color": GREEN_OK},
        {"text": "✓ Visualize trends", "size": 12, "color": GREEN_OK},
        {"text": "✓ Alert on breach", "size": 12, "color": GREEN_OK},
        {"text": "", "size": 8},
        {"text": "What They Can't Do", "bold": True, "size": 13, "color": RED_WARN},
        {"text": "✗ Predict before deployment", "size": 12, "color": RED_WARN},
        {"text": "✗ Correlate code → perf", "size": 12, "color": RED_WARN},
        {"text": "✗ Learn from history", "size": 12, "color": RED_WARN},
        {"text": "✗ CI/CD integration", "size": 12, "color": RED_WARN},
    ])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Introduction: Research Objectives
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Research Objectives")

add_textbox_lines(slide, 0.4, 0.9, 9.3, 0.5,
    [{"text": "Primary Objective:", "bold": True, "size": 16, "color": BLUE_DARK}])
add_textbox_lines(slide, 0.4, 1.2, 9.2, 0.6,
    [{"text": "Develop PerfSense — an intelligent platform that predicts performance regressions BEFORE deployment\nby analyzing code changes using machine learning.", "size": 15, "color": BLACK}])

add_rect(slide, 0.4, 1.85, 9.2, 0.04, BLUE_MID)

cols = [
    ("Build ML Model",        "> 75% accuracy\n> 70% precision\n> 60% recall",          GREEN_OK),
    ("Feature Engineering",   "28 performance-relevant features\nfrom git diffs",        BLUE_MID),
    ("Monitoring SDK",        "Core Web Vitals\n< 50 KB, < 1% overhead",                 BLUE_MID),
    ("CI/CD Integration",     "GitHub webhook\nBlock high-risk merges\n(prob ≥ 0.60)",   BLUE_MID),
]
box_w = 2.1
for i, (head, body, col) in enumerate(cols):
    x = 0.4 + i * 2.3
    add_rect(slide, x, 2.0, box_w, 2.9, LIGHT_GREY, col)
    add_textbox(slide, x+0.1, 2.05, box_w-0.2, 0.45, head,
                bold=True, size=13, color=col)
    add_rect(slide, x, 2.52, box_w, 0.03, col)
    add_textbox(slide, x+0.1, 2.6, box_w-0.2, 2.2, body,
                bold=False, size=12, color=BLACK)

add_textbox_lines(slide, 0.4, 5.1, 9.2, 0.4,
    [{"text": "Evaluate on 960 labeled commits from 8 open-source React repositories", "bold": False, "size": 13, "color": BLUE_MID, "align": PP_ALIGN.CENTER}])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Literature Review: Performance Monitoring Tools
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Literature Review: Performance Monitoring Tools")

# 5 tools: 2 rows of 2 (y=0.9, y=2.1) + 1 full-width row (y=3.3)
tools = [
    ("Lighthouse (Google)", "Automated auditing; scores LCP, FID, CLS, TTI.\nPOST-deployment only — requires a running app."),
    ("WebPageTest",         "Multi-location testing, waterfall charts.\nManual process; not suitable for CI/CD at scale."),
    ("New Relic / Datadog", "APM platforms — real-user monitoring, dashboards.\nReactive; no predictive or pre-merge capability."),
    ("React DevTools Profiler", "Component-level render timing.\nRequires running application; not integrated into CI pipeline."),
    ("Sentry",              "Error tracking + basic performance traces.  No code-change correlation or regression prediction."),
]
BOX_H = 1.1
ROW_STEP = 1.2
for i, (tool, desc) in enumerate(tools):
    if i < 4:
        row = i // 2
        col = i % 2
        x = 0.4 + col * 4.75
        y = 0.9 + row * ROW_STEP
        w = 4.4
    else:   # 5th tool spans full width
        x, y, w = 0.4, 0.9 + 2 * ROW_STEP, 9.2
    add_rect(slide, x, y, w, BOX_H, LIGHT_GREY, BLUE_MID)
    add_textbox(slide, x+0.1, y+0.05, w-0.2, 0.35, tool, bold=True, size=12.5, color=BLUE_DARK)
    add_textbox(slide, x+0.1, y+0.42, w-0.2, 0.62, desc, bold=False, size=11.5, color=BLACK)

# footer safely below all content (last row ends at 0.9+2*1.2+1.1=4.4)
add_textbox_lines(slide, 0.4, 4.55, 9.2, 0.45,
    [{"text": "Key Finding: No existing tool predicts performance impact BEFORE code is merged.", "bold": True, "size": 14, "color": RED_WARN, "align": PP_ALIGN.CENTER}])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Literature Review: ML in Software Engineering & Research Gap
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Literature Review: ML in Software Engineering")

left_lines = [
    {"text": "Defect Prediction (prior art)", "bold": True, "size": 15, "color": BLUE_DARK},
    {"text": "• XGBoost / LightGBM: 75–85% accuracy on bug prediction\n  (Chen & Guestrin, 2016; Rosen et al., 2015)", "size": 13, "color": BLACK},
    {"text": "• Domain-specific features outperform generic metrics by 15–20%", "size": 13, "color": BLACK},
    {"text": "• Class imbalance (~5–15% defect rate) addressed with SMOTE", "size": 13, "color": BLACK},
    {"text": "• Cross-project drop: 10–15% accuracy — mitigated by diverse repo selection", "size": 13, "color": BLACK},
    {"text": "", "size": 8},
    {"text": "React Performance Research", "bold": True, "size": 15, "color": BLUE_DARK},
    {"text": "• Core Web Vitals (Google 2020): LCP, FID, CLS as standard metrics", "size": 13, "color": BLACK},
    {"text": "• Unnecessary re-renders, missing memo, heavy deps = top regression causes", "size": 13, "color": BLACK},
    {"text": "• Static diff analysis proved sufficient for commit-level pattern detection", "size": 13, "color": BLACK},
]
add_textbox_lines(slide, 0.4, 0.9, 5.6, 4.4, left_lines)

# Right – research gap box
add_rect(slide, 6.2, 0.9, 3.5, 4.3, RGBColor(0xFF, 0xF2, 0xCC), RGBColor(0xFF, 0xC0, 0x00))
add_textbox_lines(slide, 6.35, 0.98, 3.2, 4.1,
    [
        {"text": "Confirmed Research Gap", "bold": True, "size": 14, "color": RGBColor(0x7F, 0x60, 0x00)},
        {"text": "", "size": 8},
        {"text": "No existing system:", "size": 13, "color": BLACK},
        {"text": "• Predicts perf regressions\n  from code diffs", "size": 12, "color": BLACK},
        {"text": "• Uses React-specific ML\n  feature engineering", "size": 12, "color": BLACK},
        {"text": "• Integrates into CI/CD\n  pipeline pre-merge", "size": 12, "color": BLACK},
        {"text": "", "size": 8},
        {"text": "PerfSense directly\naddresses this gap.", "bold": True, "size": 13, "color": BLUE_DARK},
    ])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Research Methodology: System Architecture
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Research Methodology: System Architecture",
             "6-component microservices architecture")

components = [
    ("1. Monitoring SDK",         "React 18+, TypeScript\nLCP, FID, CLS via PerformanceObserver\n< 50 KB, < 1% overhead",  "BUILT"),
    ("2. Feature Extraction",     "Python 3.11, regex diff analysis\n28 features per commit\n960 commits in ~8 min",        "BUILT"),
    ("3. ML Inference Service",   "FastAPI + XGBoost + SHAP\n/predict/diff endpoint\n< 50 ms latency",                     "BUILT"),
    ("4. Backend API Gateway",    "Node.js 18 + Express\nInput validation, UUID tracing\nPOST /api/analyze/diff",           "BUILT"),
    ("5. Web Dashboard",          "React + Recharts\nSHAP bars, risk timeline\nDark theme UI",                             "BUILT"),
    ("6. GitHub Integration",     "Webhook → backend pipeline\nBlock merge at prob ≥ 0.60\nPR status check",               "PLANNED"),
]
box_w = 2.9
box_h = 1.6
for i, (name, desc, status) in enumerate(components):
    row = i // 3
    col = i % 3
    x = 0.3 + col * 3.15
    y = 0.92 + row * 1.73   # subtitle now ends at ~0.81", safe gap before 0.92"
    color = GREEN_OK if status == "BUILT" else BLUE_MID
    add_rect(slide, x, y, box_w, box_h, LIGHT_GREY, color)
    add_textbox(slide, x+0.1, y+0.05, box_w-0.2, 0.38,
                name, bold=True, size=12, color=color)
    add_textbox(slide, x+0.1, y+0.44, box_w-0.2, 0.95,
                desc, bold=False, size=10.5, color=BLACK)
    add_textbox(slide, x+0.1, y+1.38, box_w-0.2, 0.2,
                f"[{status}]", bold=True, size=10, color=color)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Methodology: React Anti-Pattern Detection
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Research Methodology: React Anti-Pattern Framework",
             "20+ patterns across 6 categories — drives feature engineering")

cats = [
    ("Render Inefficiencies",   "Inline arrow fns, missing React.memo,\nnested component definitions",      "Weight: HIGH"),
    ("Hook Misuse",             "Unguarded useEffect (no deps array),\nexcessive useState hooks",            "Weight: HIGH"),
    ("Bundle Impact",           "Heavy dependency additions (lodash,\nmoment, chart.js), dynamic imports",   "Weight: HIGH"),
    ("Code Complexity",         "Cyclomatic complexity delta\n(for/while/if/catch additions)",               "Weight: MEDIUM"),
    ("Context Overuse",         "New Context.Provider wrappers,\nglobal state mismanagement",               "Weight: MEDIUM"),
    ("Debug Artifacts",         "console.log / console.error left in\nproduction code",                     "Weight: LOW"),
]
box_w = 4.4
box_h = 1.3   # reduced from 1.5
ROW_STEP = 1.38
for i, (cat, desc, weight) in enumerate(cats):
    row = i // 2
    col = i % 2
    x = 0.3 + col * 4.75
    y = 0.88 + row * ROW_STEP   # row2 ends at 0.88+2*1.38+1.3=4.94 — fits within 5.63"
    add_rect(slide, x, y, box_w, box_h, LIGHT_GREY, BLUE_MID)
    add_textbox(slide, x+0.1, y+0.05, box_w-0.2, 0.34,
                cat, bold=True, size=13, color=BLUE_DARK)
    add_textbox(slide, x+0.1, y+0.40, box_w-0.2, 0.66,
                desc, bold=False, size=11.5, color=BLACK)
    add_textbox(slide, x+0.1, y+1.07, box_w-0.2, 0.22,
                weight, bold=True, size=10.5, color=BLUE_MID)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Methodology: Feature Engineering (28 Features)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Research Methodology: 28-Feature Engineering",
             "Extracted from unified git diffs using Python regex — no AST required")

# Per-group: (title, body_text, box_h, body_h)
# box_h sized to fit body text: title_row=0.42" + body_h + 0.06" padding
groups = [
    ("F1–F8  Static Code Metrics",
     "lines_added, lines_removed, files_changed,\n"
     "complexity_delta, comment_delta,\n"
     "import_count_delta, test_file_changed, package_json_changed",
     1.08, 0.58),   # 3 lines ×0.19" ≈ 0.57"
    ("F9–F20  React Pattern Features",
     "arrow_fn_added, inline_fn_count, useEffect_added,\n"
     "useEffect_no_deps, useState_count, useCallback_added,\n"
     "useMemo_added, memo_added, nested_component_added,\n"
     "context_provider_added, lazy_added, console_log_added",
     1.28, 0.78),   # 4 lines ×0.19" ≈ 0.76"
    ("F21–F28  Historical & Composite",
     "F21–F23: prev_perf, trend_7d, time_since_regression"
     "  (sentinel −1: Lighthouse unavailable at scale)\n"
     "F24: large_commit  (lines_added > 200, weight +3)\n"
     "F25: heavy_dep_added  (package.json ∧ import_delta > 2)\n"
     "F26: prop_drilling_depth    F27: key_prop_missing\n"
     "F28: bundle_size_delta_proxy\n"
     "Key decision: Regex over AST — 8× faster, sufficient for diff-level counting",
     1.88, 1.38),   # 6 lines ×0.22" ≈ 1.32"
]

GAP = 0.08   # gap between boxes
y = 0.88
for title, body, box_h, body_h in groups:
    add_rect(slide, 0.3, y, 9.4, box_h, LIGHT_GREY, BLUE_DARK)
    add_textbox(slide, 0.45, y+0.05, 9.1, 0.34,
                title, bold=True, size=13, color=BLUE_DARK)
    add_rect(slide, 0.45, y+0.40, 9.1, 0.02, BLUE_DARK)   # thin separator
    add_textbox(slide, 0.45, y+0.44, 9.1, body_h,
                body, bold=False, size=11, color=BLACK)
    y += box_h + GAP
# groups end at ~0.88+1.08+0.08+1.28+0.08+1.88=5.28" — no footer needed


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Data Collection: Repository Selection
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Data Collection: Repository Selection",
             "8 open-source React repositories — 120 non-merge commits each = 960 total")

repos = [
    ("facebook/react",         "Core library",    "Small",  "120", "12.0%"),
    ("excalidraw/excalidraw",  "Drawing tool",    "Medium", "120", "11.7%"),
    ("ant-design/ant-design-pro", "Enterprise UI","Large",  "120",  "8.3%"),
    ("facebook/docusaurus",    "Docs site",       "Medium", "120", "10.8%"),
    ("discourse/discourse",    "Forum platform",  "Large",  "120",  "1.7%"),
    ("strapi/strapi",          "CMS backend",     "Large",  "120", "29.2%"),
    ("laurent22/joplin",       "Note-taking app", "Medium", "120", "18.3%"),
    ("calcom/cal.com",         "Calendar SaaS",   "Large",  "120", "27.5%"),
]

# header row
hdr_y = 0.9
add_rect(slide, 0.3, hdr_y, 9.4, 0.38, BLUE_DARK)
for col_x, lbl in [(0.4, "Repository"), (3.5, "Category"), (5.3, "Size"),
                   (6.5, "Commits"), (7.9, "Regression %")]:
    add_textbox(slide, col_x, hdr_y+0.04, 1.6, 0.3, lbl,
                bold=True, size=11, color=WHITE)

for i, (repo, cat, size, commits, pct) in enumerate(repos):
    y = 1.3 + i * 0.47
    fill = LIGHT_GREY if i % 2 == 0 else WHITE
    add_rect(slide, 0.3, y, 9.4, 0.44, fill)
    for col_x, txt in [(0.4, repo), (3.5, cat), (5.3, size),
                       (6.5, commits), (7.9, pct)]:
        add_textbox(slide, col_x, y+0.06, 1.6, 0.32,
                    txt, bold=False, size=11, color=BLACK)

add_textbox_lines(slide, 0.3, 5.08, 9.4, 0.42,
    [{"text": "Overall: 141 / 960 regressions (14.7%) — within the expected 10–15% class imbalance range.", "bold": True, "size": 13, "color": BLUE_DARK}])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Data Collection: Heuristic Proxy Labeling
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Data Collection: Heuristic Proxy Labeling",
             "Lighthouse CI on 960 historical commits ≈ 48 hrs compute — not feasible")

add_textbox_lines(slide, 0.4, 0.92, 5.5, 1.95,
    [
        {"text": "Labeling Strategy", "bold": True, "size": 15, "color": BLUE_DARK},
        {"text": "Each commit scored on 16 signals:", "size": 13, "color": BLACK},
        {"text": "  +3  large_commit  (lines_added > 200)", "size": 12, "color": RED_WARN},
        {"text": "  +2  heavy_dep_added, 9 other positive signals", "size": 12, "color": RED_WARN},
        {"text": "  −2  useCallback_added, 4 protective signals", "size": 12, "color": GREEN_OK},
        {"text": "  Auto-tuned threshold = 3  →  score ≥ threshold = regression", "size": 12, "color": BLUE_MID},
    ])

add_textbox_lines(slide, 0.4, 2.95, 5.5, 2.0,
    [
        {"text": "Lighthouse Validation (10 commits)", "bold": True, "size": 15, "color": BLUE_DARK},
        {"text": "• All 10 commits scored 51–55 (stable)", "size": 13, "color": BLACK},
        {"text": "• 6/10 agreement with heuristic labels", "size": 13, "color": BLACK},
        {"text": "• 4 false positives (large but stable commits)", "size": 13, "color": RED_WARN},
        {"text": "• 0 false negatives — conservative bias ✓", "size": 13, "color": GREEN_OK},
    ])

# Right side stats — height 3.9 so ends at 0.92+3.9=4.82, footer at 4.95 safe
add_rect(slide, 6.1, 0.92, 3.6, 3.9, LIGHT_GREY, BLUE_DARK)
add_textbox_lines(slide, 6.25, 1.0, 3.3, 3.7,
    [
        {"text": "Dataset Statistics", "bold": True, "size": 14, "color": BLUE_DARK},
        {"text": "", "size": 7},
        {"text": "Total commits:    960", "size": 13, "color": BLACK},
        {"text": "Regressions:      141", "size": 13, "color": RED_WARN, "bold": True},
        {"text": "Clean commits:    819", "size": 13, "color": GREEN_OK},
        {"text": "Regression rate:  14.7%", "size": 13, "color": BLACK},
        {"text": "", "size": 7},
        {"text": "Train / Val / Test Split", "bold": True, "size": 13, "color": BLUE_DARK},
        {"text": "70% / 10% / 20%  (chronological)", "size": 12, "color": BLACK},
        {"text": "SMOTE on training set only", "size": 12, "color": BLACK},
        {"text": "Saved: data/features/features_labeled.csv", "size": 11, "color": BLUE_MID},
    ])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – ML Results
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "ML Model Training & Results",
             "Three models trained; XGBoost selected as primary")

# Model comparison table
add_rect(slide, 0.3, 0.9, 9.4, 0.4, BLUE_DARK)
for col_x, lbl in [(0.4, "Model"), (3.2, "Accuracy"), (5.0, "Precision"), (6.8, "Recall"), (8.3, "F1"), (9.2, "AUC")]:
    add_textbox(slide, col_x, 0.94, 1.8, 0.32, lbl, bold=True, size=12, color=WHITE)

models = [
    ("Logistic Regression", "0.912", "0.748", "0.846", "0.794", "0.966",  False),
    ("Random Forest",       "0.959", "0.883", "0.885", "0.884", "0.983",  False),
    ("XGBoost ★ PRIMARY",   "0.974", "0.938", "0.900", "0.918", "0.9883", True),
]
for i, (name, acc, prec, rec, f1, auc, primary) in enumerate(models):
    y = 1.32 + i * 0.52
    fill = RGBColor(0xE2, 0xEF, 0xDA) if primary else (LIGHT_GREY if i % 2 == 0 else WHITE)
    add_rect(slide, 0.3, y, 9.4, 0.5, fill)
    col = GREEN_OK if primary else BLACK
    bold = primary
    for col_x, txt in [(0.4, name), (3.2, acc), (5.0, prec), (6.8, rec), (8.3, f1), (9.2, auc)]:
        add_textbox(slide, col_x, y+0.07, 1.8, 0.36,
                    txt, bold=bold, size=12, color=col)

# Targets met box
add_rect(slide, 0.3, 3.0, 4.5, 2.4, RGBColor(0xE2, 0xEF, 0xDA), GREEN_OK)
add_textbox_lines(slide, 0.45, 3.08, 4.2, 2.3,
    [
        {"text": "Dissertation Targets — ALL MET", "bold": True, "size": 13, "color": GREEN_OK},
        {"text": "✓  Accuracy  > 75%  →  97.4%", "size": 12, "color": GREEN_OK},
        {"text": "✓  Precision > 70%  →  93.8%", "size": 12, "color": GREEN_OK},
        {"text": "✓  Recall    > 60%  →  90.0%", "size": 12, "color": GREEN_OK},
        {"text": "✓  F1-Score        →  91.8%", "size": 12, "color": GREEN_OK},
        {"text": "✓  AUC-ROC         →  0.9883", "size": 12, "color": GREEN_OK},
    ])

# Config box
add_rect(slide, 5.0, 3.0, 4.7, 2.4, LIGHT_GREY, BLUE_MID)
add_textbox_lines(slide, 5.15, 3.08, 4.4, 2.3,
    [
        {"text": "XGBoost Configuration", "bold": True, "size": 13, "color": BLUE_DARK},
        {"text": "n_estimators: 500  |  max_depth: 5", "size": 11.5, "color": BLACK},
        {"text": "learning_rate: 0.05  |  subsample: 0.8", "size": 11.5, "color": BLACK},
        {"text": "colsample_bytree: 0.8", "size": 11.5, "color": BLACK},
        {"text": "Grid search: 5-fold CV on val set", "size": 11.5, "color": BLACK},
        {"text": "SHAP TreeExplainer for interpretability", "size": 11.5, "color": BLACK},
    ])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – SHAP Feature Importance
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "SHAP Feature Importance",
             "SHAP TreeExplainer — mean |SHAP| across test set")

features = [
    ("lines_added",            0.412, "Commit size — raw lines added"),
    ("large_commit",           0.371, "Binary: lines_added > 200 (weight +3)"),
    ("complexity_delta",       0.289, "Cyclomatic complexity added (for/if/while)"),
    ("useEffect_added",        0.187, "New useEffect hooks in the commit"),
    ("heavy_dep_added",        0.156, "Heavy dependency + import growth"),
    ("useEffect_no_deps",      0.134, "useEffect without dependency array"),
    ("nested_component_added", 0.112, "Component defined inside another component"),
    ("import_count_delta",     0.098, "Net new import statements"),
    ("inline_fn_count",        0.087, "Inline arrow functions as props"),
    ("files_changed",          0.074, "Number of files touched"),
]

max_val = 0.412
bar_area_w = 5.0
bar_start_x = 3.7

for i, (fname, val, note) in enumerate(features):
    y = 0.9 + i * 0.45
    # label
    add_textbox(slide, 0.3, y, 3.3, 0.4, fname, bold=(i < 3), size=11.5,
                color=BLUE_DARK if i < 3 else BLACK)
    # bar
    bar_w = bar_area_w * val / max_val
    bar_color = BLUE_DARK if i == 0 else (BLUE_MID if i < 3 else RGBColor(0x9D, 0xC3, 0xE6))
    add_rect(slide, bar_start_x, y+0.06, bar_w, 0.28, bar_color)
    # value
    add_textbox(slide, bar_start_x + bar_w + 0.05, y, 0.5, 0.4,
                f"{val:.3f}", bold=(i < 3), size=10, color=BLACK)
    # note
    add_textbox(slide, bar_start_x + bar_w + 0.62, y, 2.8, 0.4,
                note, bold=False, size=9.5, color=RGBColor(0x60, 0x60, 0x60))

add_textbox_lines(slide, 0.3, 5.08, 9.4, 0.42,
    [{"text": "Insight: Commit size (lines_added, large_commit) dominates. React-specific features add signal on top.", "bold": True, "size": 12, "color": BLUE_MID}])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – Cross-Repository Validation
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Cross-Repository Validation (Leave-One-Repo-Out)",
             "Model trained on 7 repos, tested on held-out repo — 8 runs")

# Header row
add_rect(slide, 0.3, 0.9, 9.4, 0.38, BLUE_DARK)
for col_x, lbl in [(0.4, "Left-Out Repo"), (3.4, "F1"), (5.0, "AUC"), (6.5, "Precision"), (8.0, "Recall")]:
    add_textbox(slide, col_x, 0.94, 2.5, 0.3, lbl, bold=True, size=12, color=WHITE)

cross_repos = [
    ("ant-design-pro", "1.000", "0.998", "1.000", "1.000"),
    ("calcom/cal.com",  "0.900", "0.991", "0.900", "0.900"),
    ("discourse",       "1.000", "0.999", "1.000", "1.000"),
    ("docusaurus",      "0.923", "0.995", "0.923", "0.923"),
    ("facebook/react",  "0.895", "0.982", "0.909", "0.882"),
    ("joplin",          "0.833", "0.982", "0.857", "0.812"),
    ("strapi",          "0.941", "0.998", "0.941", "0.941"),
    ("excalidraw",      "0.917", "0.991", "0.917", "0.917"),
]
for i, (repo, f1, auc, prec, rec) in enumerate(cross_repos):
    y = 1.3 + i * 0.46
    fill = LIGHT_GREY if i % 2 == 0 else WHITE
    add_rect(slide, 0.3, y, 9.4, 0.44, fill)
    f1_color = GREEN_OK if float(f1) >= 0.9 else (RGBColor(0xFF, 0x70, 0x00) if float(f1) >= 0.85 else RED_WARN)
    add_textbox(slide, 0.4, y+0.06, 2.8, 0.32, repo, bold=False, size=12, color=BLACK)
    add_textbox(slide, 3.4, y+0.06, 1.4, 0.32, f1,   bold=True,  size=12, color=f1_color)
    add_textbox(slide, 5.0, y+0.06, 1.4, 0.32, auc,  bold=False, size=12, color=BLACK)
    add_textbox(slide, 6.5, y+0.06, 1.4, 0.32, prec, bold=False, size=12, color=BLACK)
    add_textbox(slide, 8.0, y+0.06, 1.4, 0.32, rec,  bold=False, size=12, color=BLACK)

add_textbox_lines(slide, 0.3, 5.08, 9.4, 0.42,
    [{"text": "F1 range: 0.833 – 1.000 across all repos. AUC consistently > 0.98. Confirms model generalizes beyond training data.", "bold": True, "size": 13, "color": GREEN_OK}])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – Case Studies
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Real-World Validation: Case Studies",
             "3 synthetic React apps — deliberate regressions — end-to-end prediction pipeline")

cases = [
    (
        "CS1: Task Manager",
        "58.8% — MEDIUM",
        BLUE_MID,
        "4 unguarded useEffects, 2 nested\ncomponents, 8+ inline arrows,\nremoved all memoization\n(69-line diff → large_commit=0)",
        "Correct: 69 lines, large_commit\nnot triggered. Model predicts\nMEDIUM risk accurately.",
    ),
    (
        "CS2: Product Catalog",
        "81.8% — HIGH",
        RED_WARN,
        "lodash/moment/chart.js/xlsx added,\n152+ lines, 5 unguarded useEffects,\nnested ProductCard + StatsBar",
        "Correct: heavy_dep_added\nfired + large_commit. HIGH risk\ncorrectly predicted.",
    ),
    (
        "CS3: Analytics Widget",
        "98.9% — HIGH",
        RED_WARN,
        "complexity_delta=12, 9 useState,\n3 unguarded useEffects,\n2 nested components,\n11 console.logs",
        "Correct: highest complexity\ndelta in dataset. Extreme HIGH\nrisk — model confident.",
    ),
]

for i, (title, result, res_color, code_desc, analysis) in enumerate(cases):
    x = 0.3 + i * 3.2
    add_rect(slide, x, 0.88, 3.0, 3.85, LIGHT_GREY, BLUE_DARK)  # height 3.85 → ends at 4.73
    add_textbox(slide, x+0.08, 0.95, 2.84, 0.40, title, bold=True, size=13, color=BLUE_DARK)
    add_rect(slide, x+0.08, 1.38, 2.84, 0.36, res_color)
    add_textbox(slide, x+0.08, 1.40, 2.84, 0.34, result, bold=True, size=12.5, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x+0.08, 1.82, 2.84, 1.55, code_desc, bold=False, size=11, color=BLACK)
    add_rect(slide, x+0.08, 3.38, 2.84, 0.03, BLUE_MID)
    add_textbox(slide, x+0.08, 3.44, 2.84, 1.0, analysis, bold=False, size=11, color=BLUE_MID)

# footer well below content (boxes end at 0.88+3.85=4.73)
add_textbox_lines(slide, 0.3, 4.85, 9.4, 0.42,
    [{"text": "Result: 3/3 regressions correctly detected. Zero false negatives. End-to-end latency < 200 ms.", "bold": True, "size": 13, "color": GREEN_OK}])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 – Implementation Status
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Implementation Status",
             "Ahead of original schedule — ML + Services + Dashboard completed in Month 3")

tasks = [
    # (phase, task, status, notes)
    ("Month 1", "Literature Review",           "DONE", "20+ papers; gap confirmed"),
    ("Month 1", "Architecture Design",         "DONE", "6-component microservices"),
    ("Month 1", "Repository Selection",        "DONE", "8 repos selected"),
    ("Month 2", "Commit Extraction Pipeline",  "DONE", "960 commits, 8 repos"),
    ("Month 2", "Feature Engineering (28F)",   "DONE", "Python regex diff analysis"),
    ("Month 2", "Heuristic Proxy Labeling",    "DONE", "141/960 regressions (14.7%)"),
    ("Month 3", "ML Training + Grid Search",   "DONE", "F1=0.918, AUC=0.9883"),
    ("Month 3", "FastAPI Inference Service",   "DONE", "/predict/diff endpoint live"),
    ("Month 3", "Node.js/Express Backend",     "DONE", "Input validation, UUID tracing"),
    ("Month 3", "Monitoring SDK (TS)",         "DONE", "LCP/FID/CLS + React Profiler"),
    ("Month 3", "React Dashboard",             "DONE", "SHAP bars, risk timeline"),
    ("Month 3", "3 Case Studies",              "DONE", "3/3 regressions detected"),
    ("Month 4", "GitHub PR Integration",       "TODO", "Webhook + status check"),
    ("Month 4", "Dissertation Writing",        "TODO", "80–100 pages"),
]

add_rect(slide, 0.3, 0.9, 9.4, 0.38, BLUE_DARK)
for col_x, lbl in [(0.4, "Phase"), (1.55, "Task"), (5.8, "Status"), (7.1, "Notes")]:
    add_textbox(slide, col_x, 0.94, 2.5, 0.3, lbl, bold=True, size=11, color=WHITE)

for i, (phase, task, status, notes) in enumerate(tasks):
    y = 1.3 + i * 0.29
    fill = LIGHT_GREY if i % 2 == 0 else WHITE
    add_rect(slide, 0.3, y, 9.4, 0.27, fill)
    s_color = GREEN_OK if status == "DONE" else BLUE_MID
    add_textbox(slide, 0.4,  y+0.03, 1.1,  0.22, phase,  bold=False, size=9.5, color=BLACK)
    add_textbox(slide, 1.55, y+0.03, 4.2,  0.22, task,   bold=False, size=9.5, color=BLACK)
    add_textbox(slide, 5.8,  y+0.03, 1.2,  0.22, status, bold=True,  size=9.5, color=s_color)
    add_textbox(slide, 7.1,  y+0.03, 2.55, 0.22, notes,  bold=False, size=9.0, color=RGBColor(0x40, 0x40, 0x40))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 – Challenges & Mitigation
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Challenges & Mitigation Strategies")

challenges = [
    (
        "Lighthouse CI Scale (960 commits ≈ 48 hrs)",
        "Heuristic proxy labeling (16 signals, auto-tuned threshold=3).\n"
        "Validated on 10-commit Lighthouse sample: 0 false negatives.",
        RED_WARN,
    ),
    (
        "Class Imbalance (14.7% regression rate)",
        "SMOTE oversampling on training set only.\n"
        "Preserved class distribution in validation and test sets.",
        RED_WARN,
    ),
    (
        "Cross-Project Generalization",
        "8 diverse repos (small→large, tool→platform).\n"
        "Leave-one-repo-out CV confirms F1 ≥ 0.833 for every repo.",
        RED_WARN,
    ),
    (
        "Python Regex Lookbehind Constraints",
        "Fixed-width lookbehind in re module — patterns simplified.\n"
        "e.g. r'React\\.memo\\s*\\(|=\\s*memo\\s*\\(' for memo detection.",
        BLUE_MID,
    ),
    (
        "Windows Unicode / Port Management",
        "sys.stdout.reconfigure(encoding='utf-8') throughout.\n"
        "netstat + taskkill replaces npx kill-port (which timed out).",
        BLUE_MID,
    ),
]

for i, (challenge, mitigation, col) in enumerate(challenges):
    y = 0.9 + i * 0.87
    add_rect(slide, 0.3, y, 9.4, 0.82, LIGHT_GREY, col)
    add_textbox(slide, 0.42, y+0.04, 9.1, 0.32, f"Challenge: {challenge}",
                bold=True, size=12, color=col)
    add_textbox(slide, 0.42, y+0.38, 9.1, 0.4,  f"Mitigation: {mitigation}",
                bold=False, size=11.5, color=BLACK)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 – Next Steps & Timeline
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Next Steps & Revised Timeline",
             "On track for June 2026 submission")

months = [
    ("Month 1\n(Feb 2026)",  "DONE", GREEN_OK,
     "Literature review\nArchitecture design\nRepo selection"),
    ("Month 2\n(Mar 2026)",  "DONE", GREEN_OK,
     "Data collection (960)\nFeature engineering\nProxy labeling"),
    ("Month 3\n(Apr 2026)",  "DONE", GREEN_OK,
     "ML training (F1=0.918)\nFastAPI service\nBackend + Dashboard\nCase studies"),
    ("Month 4\n(May–Jun 2026)", "IN PROGRESS", BLUE_MID,
     "GitHub PR integration\nFinal evaluation\nAblation study\nDissertation writing\n(80–100 pages)"),
]

for i, (month, status, col, tasks_text) in enumerate(months):
    x = 0.3 + i * 2.4
    add_rect(slide, x, 0.88, 2.1, 3.85, LIGHT_GREY, col)  # height 3.85 → ends at 4.73
    add_textbox(slide, x+0.08, 0.93, 1.94, 0.52, month, bold=True, size=12.5, color=col)
    add_rect(slide, x+0.08, 1.48, 1.94, 0.32, col)
    add_textbox(slide, x+0.08, 1.50, 1.94, 0.28, status, bold=True, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x+0.08, 1.86, 1.94, 2.75, tasks_text, bold=False, size=11, color=BLACK)

# footer well below content (boxes end at 0.88+3.85=4.73)
add_textbox_lines(slide, 0.3, 4.85, 9.4, 0.42,
    [{"text": "Immediate priority: GitHub PR webhook integration → dissertation writing (targeting final submission by June 2026)", "size": 13, "color": BLUE_MID}])


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 – Summary
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
slide_header(slide, "Summary")

points = [
    ("Problem",        "Performance regressions discovered post-deployment — costing revenue and developer time.",         BLUE_DARK),
    ("Gap",            "No existing tool predicts React performance impact from code changes before merge.",                RED_WARN),
    ("Solution",       "PerfSense: ML pipeline that analyzes git diffs, extracts 28 features, and predicts regressions.",  BLUE_MID),
    ("Data",           "960 commits from 8 open-source React repos; 141 regressions (14.7%); chronological split.",        BLACK),
    ("Model",          "XGBoost: F1=0.918, AUC=0.9883, Accuracy=97.4% — all dissertation targets exceeded.",              GREEN_OK),
    ("Generalization", "Leave-one-repo-out: F1 ≥ 0.833 for every repo. AUC > 0.98 consistently.",                         GREEN_OK),
    ("Validation",     "3/3 real-world case studies correctly classified; end-to-end latency < 200 ms.",                   GREEN_OK),
    ("Status",         "Months 1–3 COMPLETE (ahead of schedule). Month 4: GitHub integration + dissertation writing.",     BLUE_MID),
]

for i, (label, text, col) in enumerate(points):
    y = 0.92 + i * 0.55
    add_rect(slide, 0.3, y, 1.6, 0.44, col)
    add_textbox(slide, 0.32, y+0.04, 1.56, 0.36, label, bold=True, size=12, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 2.1, y+0.05, 7.7, 0.44, text, bold=False, size=12.5, color=col if col != BLACK else RGBColor(0x20, 0x20, 0x20))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20 – Thank You
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = WHITE

add_rect(slide, 0, 0, 10, 5.63, WHITE)
add_textbox(slide, 0.5, 1.0, 9.0, 1.0,
            "Thank You",
            bold=True, size=44, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_textbox(slide, 0.5, 2.1, 9.0, 0.5,
            "Questions & Discussion",
            bold=False, size=24, color=BLUE_MID, align=PP_ALIGN.CENTER)
add_rect(slide, 2.0, 2.75, 6.0, 0.05, BLUE_MID)
add_textbox_lines(slide, 0.5, 2.88, 9.0, 2.3,
    [
        {"text": "CHETHAN S",                               "bold": True,  "size": 16, "color": BLACK, "align": PP_ALIGN.CENTER},
        {"text": "BITS ID: 2024TM93208",                    "size": 14,   "color": BLACK, "align": PP_ALIGN.CENTER},
        {"text": "MTech Software Engineering — BITS Pilani WILP", "size": 13, "color": BLACK, "align": PP_ALIGN.CENTER},
        {"text": "Supervisor: Ankit Kagliwal",              "size": 13,   "color": BLACK, "align": PP_ALIGN.CENTER},
        {"text": "chethans@cadence.com",                    "size": 13,   "color": BLUE_MID, "align": PP_ALIGN.CENTER},
    ])


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "docs/reports/MS2024TM93208_MidSem_Viva.pptx"
prs.save(out_path)
print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
